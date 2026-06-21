# -*- coding: utf-8 -*-
"""補 3 張 Canvas(黑金風,配合 RWA 0607 deck):Tokenomics / Use Case / Smart Contract Role Map。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG    = RGBColor(0x0A, 0x0A, 0x0B)
GOLD  = RGBColor(0xF5, 0xB5, 0x44)
ORANGE= RGBColor(0xE8, 0x55, 0x2D)
WHITE = RGBColor(0xF2, 0xF2, 0xF2)
MUT   = RGBColor(0xB4, 0xB4, 0xB8)
CARD  = RGBColor(0x17, 0x17, 0x1A)
LINE  = RGBColor(0x33, 0x30, 0x2A)
HEAD  = RGBColor(0x24, 0x1c, 0x10)
FONT  = "Microsoft JhengHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    # 底部金色條
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()
    # 左上：金色神殿 logo(簡化:屋頂 + 底座 + 兩根柱)
    roof = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.5), Inches(0.4), Inches(0.5), Inches(0.18))
    roof.fill.solid(); roof.fill.fore_color.rgb = GOLD; roof.line.fill.background()
    base = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.78), Inches(0.5), Inches(0.08))
    base.fill.solid(); base.fill.fore_color.rgb = GOLD; base.line.fill.background()
    for dx in (0.56, 0.71, 0.86):
        c = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(dx), Inches(0.6), Inches(0.05), Inches(0.18))
        c.fill.solid(); c.fill.fore_color.rgb = GOLD; c.line.fill.background()
    return s


def tb(s, l, t, w, h):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); b.text_frame.word_wrap = True
    return b.text_frame


def setp(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, space=6):
    p.text = text; p.alignment = align; p.space_after = Pt(space)
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT


def title(s, text):
    tf = tb(s, 1.25, 0.42, 11.5, 0.9)
    setp(tf.paragraphs[0], text, 32, WHITE, bold=True)


def rows(s, items, top=1.7, size=15, gap=9, label_color=GOLD):
    """items: (label, text) -> 「▍ label：text」label 上色、text 白。"""
    tf = tb(s, 0.7, top, 12.0, 7.5 - top - 0.5)
    for i, (label, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▍ " + label + "  "
        p.runs[0].font.size = Pt(size); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = label_color; p.runs[0].font.name = FONT
        r2 = p.add_run(); r2.text = text
        r2.font.size = Pt(size); r2.font.color.rgb = WHITE; r2.font.name = FONT
        p.space_after = Pt(gap)


def table(s, headers, data, top=1.7, left=0.7, width=12.0, fsize=14):
    nr = len(data) + 1
    t = s.shapes.add_table(nr, len(headers), Inches(left), Inches(top), Inches(width), Inches(0.42 * nr)).table
    for ci, h in enumerate(headers):
        c = t.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = HEAD; c.text = h
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        if c.text_frame.paragraphs[0].runs:
            rr = c.text_frame.paragraphs[0].runs[0]; rr.font.size = Pt(fsize); rr.font.bold = True; rr.font.color.rgb = GOLD; rr.font.name = FONT
    for r, row in enumerate(data):
        for ci, val in enumerate(row):
            c = t.cell(r + 1, ci); c.fill.solid(); c.fill.fore_color.rgb = CARD; c.text = val
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            if c.text_frame.paragraphs[0].runs:
                rr = c.text_frame.paragraphs[0].runs[0]; rr.font.size = Pt(fsize); rr.font.name = FONT
                rr.font.color.rgb = GOLD if ci == 0 else WHITE


# ---------- 1 Tokenomics Canvas ----------
s = slide(); title(s, "Tokenomics Canvas ｜ 代幣經濟")
rows(s, [
    ("代幣名稱與定位", "dXXXX（如 dTSMC）= 代幣化台股,1 代幣 ↔ 1 股經濟權利;TWD = 平台計價穩定幣。"),
    ("背後資產與來源", "台灣上市股票;價格來源 TWSE（預言機餵價）;商轉以實股 1:1 託管擔保。"),
    ("鑄造 / 銷毀 Mint-Burn", "存 TWD 依預言機即時價 mint;redeem 反向燒毀;穩定幣 mintTWD / requestRedemption。"),
    ("發行量與增發規則", "無固定上限,供給隨「TWD 抵押鑄造」彈性增減,恆等於儲備擔保量。"),
    ("權限角色", "admin / issuer · ATTESTOR（儲備簽證）· COMPLIANCE（凍結 / 暫停）· 工廠 owner。"),
    ("費用結構", "Demo 免交易費;商轉開鑄造 / 贖回費、AUM 管理費、借貸利差、IPO 承銷費。"),
    ("二級市場與流動性", "平台 mint / redeem + Uniswap v3 池（已上鏈）+ 借貸釋放抵押流動性。"),
], top=1.75, size=15, gap=10)

# ---------- 2 Use Case Canvas ----------
s = slide(); title(s, "Use Case Canvas ｜ 使用情境")
rows(s, [
    ("目標使用者", "機構（發行 / 做市）、個人（配置 / 賺息）、零售（零股 / 跨境）。"),
    ("使用場景與痛點解方", "夜間配置台股、零股碎片化、抵押借 TWD、SpaceX / 新股認購、即時驗證估值與擔保。"),
    ("鏈上 / 鏈下流程", "鏈下 TWSE 餵價 + AI 估值 → 鏈上 mint / redeem / 借貸 / 認購 / Uniswap;前端 wagmi 直接讀寫。"),
    ("法規與合規挑戰", "證券型代幣定性、KYC / AML、穩定幣儲備與贖回、跨境稅務。"),
    ("競品分析與差異化", "vs 券商（24/7+碎片化）、vs CEX（鏈上可驗擔保）、vs Ondo / xStocks（專注台股 + TWD 本位）。"),
    ("風險管理機制", "凍結 setFrozen、暫停 pause、黑名單、預言機過期價檢查、借貸健康因子 / 清算。"),
], top=1.75, size=15, gap=11, label_color=ORANGE)

# ---------- 3 Smart Contract Role Map ----------
s = slide(); title(s, "Smart Contract Role Map ｜ 合約角色")
table(s, ["角色", "權限"], [
    ["DEFAULT_ADMIN / issuer", "發行、設定揭露、issue 鑄造"],
    ["ATTESTOR", "attestReserves 簽署儲備證明（時間戳 + 文件雜湊）"],
    ["COMPLIANCE", "setFrozen 凍結地址、pause / unpause 全面暫停"],
    ["owner（工廠 / 借貸 / IPO）", "上架代幣、設風險參數、建立 IPO 認購案"],
    ["feeder", "updatePrices 餵價"],
], top=1.7, fsize=14)
rows(s, [
    ("系統合約", "RegulatedTWD · PriceOracle · StockFactory · StockToken×12 · LendingPool · StockIPO · Uniswap v3 池（皆驗證）。"),
    ("可升級性", "現為不可升級（易審計、降信任風險）;路線圖評估 UUPS Proxy + 多簽 / 時間鎖。"),
    ("資產轉移邏輯", "_update 覆寫攔截黑名單 / 凍結 from-to;whenNotPaused 守鑄造 / 贖回。"),
    ("簽名轉移 permit", "路線圖支援 ERC-2612（免 gas 授權,優化 UX）。"),
    ("事件與透明性", "Minted / Redeemed / PriceUpdated / ReserveAttested / AddressFrozen / Borrow / Liquidate 全上鏈可稽核。"),
], top=4.55, size=12.5, gap=6)

out = r"C:\Users\lab643\Desktop\rwa report\FormosaX_3Canvas_補充.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
