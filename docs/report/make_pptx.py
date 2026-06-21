# -*- coding: utf-8 -*-
"""生成 FormosaX RWA pitch 的 PowerPoint(深色玉金主題)。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG    = RGBColor(0x07, 0x09, 0x0E)
CARD  = RGBColor(0x0E, 0x14, 0x1F)
JADE  = RGBColor(0x19, 0xE6, 0xB0)
GOLD  = RGBColor(0xF5, 0xB5, 0x44)
FG    = RGBColor(0xE7, 0xEE, 0xF7)
MUT   = RGBColor(0x90, 0x9D, 0xAE)
LINE  = RGBColor(0x26, 0x30, 0x3E)
HEADBG= RGBColor(0x12, 0x1A, 0x26)
FONT  = "Microsoft JhengHei"
MONO  = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    # 左上品牌條
    bar = s.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.32), Inches(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = JADE
    bar.line.fill.background()
    bn = s.shapes.add_textbox(Inches(1.0), Inches(0.43), Inches(3), Inches(0.4)).text_frame
    bn.text = "FormosaX"
    r = bn.paragraphs[0].runs[0]; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = MUT; r.font.name = FONT
    return s


def tb(s, l, t, w, h):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    return tf


def setp(p, text, size, color, bold=False, font=FONT, align=PP_ALIGN.LEFT, space=6):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space)
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = font
    return p


def kicker(s, text):
    tf = tb(s, 0.9, 1.05, 11, 0.5)
    setp(tf.paragraphs[0], text, 13, JADE, bold=True, font=MONO)


def title(s, text, size=34, color=FG):
    tf = tb(s, 0.9, 1.45, 11.5, 1.2)
    setp(tf.paragraphs[0], text, size, color, bold=True)


def bullets(s, items, top=2.7, size=16, width=11.5, gap=8):
    tf = tb(s, 0.9, top, width, 7.5 - top - 0.4)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # 支援 (粗體前綴, 說明)
        if isinstance(it, tuple):
            head, rest = it
            p.text = "▍ " + head + "  "
            p.runs[0].font.size = Pt(size); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = JADE; p.runs[0].font.name = FONT
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = FG; r2.font.name = FONT
        else:
            p.text = "▍ " + it
            p.runs[0].font.size = Pt(size); p.runs[0].font.color.rgb = FG; p.runs[0].font.name = FONT
        p.space_after = Pt(gap)


TOTAL = 17
def pageno(s, n):
    tf = tb(s, 11.6, 6.95, 1.5, 0.4)
    setp(tf.paragraphs[0], f"{n:02d} / {TOTAL}", 11, MUT, font=MONO, align=PP_ALIGN.RIGHT)


def table(s, headers, rows, top=2.7, left=0.9, width=11.5, fsize=13):
    nr, nc = len(rows) + (1 if headers else 0), len(rows[0])
    gtbl = s.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(width), Inches(0.4 * nr)).table
    ri = 0
    if headers:
        for ci, h in enumerate(headers):
            c = gtbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = HEADBG
            c.text = h; pr = c.text_frame.paragraphs[0].runs
            if pr:
                pr[0].font.size = Pt(fsize); pr[0].font.bold = True; pr[0].font.color.rgb = JADE; pr[0].font.name = FONT
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
        ri = 1
    for r, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = gtbl.cell(r + ri, ci); c.fill.solid(); c.fill.fore_color.rgb = CARD
            c.text = val
            pr = c.text_frame.paragraphs[0].runs
            if pr:
                pr[0].font.size = Pt(fsize); pr[0].font.name = FONT
                pr[0].font.color.rgb = JADE if ci == nc - 1 and headers else FG
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
    return gtbl


def cards(s, items, top=2.7, cols=3):
    gap = 0.3
    total = 11.5
    w = (total - gap * (cols - 1)) / cols
    for i, (h, body, accent) in enumerate(items):
        col = i % cols; row = i // cols
        l = 0.9 + col * (w + gap); t = top + row * 2.0
        box = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(1.8))
        box.fill.solid(); box.fill.fore_color.rgb = CARD
        box.line.color.rgb = accent if accent else LINE; box.line.width = Pt(1)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.14)
        setp(tf.paragraphs[0], h, 15, accent if accent else JADE, bold=True)
        p = tf.add_paragraph(); setp(p, body, 12, MUT)


# ---------- 1 封面 ----------
s = slide()
mk = s.shapes.add_shape(1, Inches(0.9), Inches(1.6), Inches(1.0), Inches(1.0))
mk.fill.solid(); mk.fill.fore_color.rgb = JADE; mk.line.fill.background()
tf = tb(s, 0.9, 0.95, 11, 0.5); setp(tf.paragraphs[0], "REAL WORLD ASSET · ON-CHAIN · AI VALUATION", 13, JADE, bold=True, font=MONO)
tf = tb(s, 0.9, 2.85, 11.8, 2.0)
setp(tf.paragraphs[0], "FormosaX — 把台股變成 24/7 可碎片化的鏈上資產", 40, FG, bold=True)
tf = tb(s, 0.9, 4.45, 11.5, 1.2)
setp(tf.paragraphs[0], "受監管 TWD 穩定幣 × 預言機即時估值 × AI 輔助公允估值/風險 × 鏈上 1:1 實股擔保 × 可組合 DeFi(借貸/IPO)", 16, MUT)
tf = tb(s, 0.9, 5.7, 11.8, 1.3)
setp(tf.paragraphs[0], "題目七 · RWA 資產代幣化與估值系統   |   組員 114AB8049 · 114AB8031 · 114AB8051 · 114AB8016", 13, FG, font=MONO)
setp(tf.add_paragraph(), "Live Demo:https://rwa0607.vercel.app   |   Ethereum Sepolia(Etherscan 已驗證)", 13, JADE, font=MONO)
pageno(s, 1)

# ---------- 2 商業情境 ----------
s = slide(); kicker(s, "BUSINESS CONTEXT"); title(s, "商業情境")
bullets(s, [
    "金融機構希望把不動產、債權、股票等資產上鏈,提高流動性與投資可及性 —— 但卡在兩道牆:",
], top=2.55, size=16)
cards(s, [
    ("① 估值不透明", "資產缺乏即時、可驗證的公允價格,投資人難判斷標的價值與背後擔保。", GOLD),
    ("② 流動性受限", "交易受時段、單位(整股)、跨境與 T+2 結算限制,資金進出摩擦高。", JADE),
    ("FormosaX 切入", "以台股切入(資料公開、估值明確、最易驗證),示範可複製的代幣化×預言機×AI×DeFi 堆疊。", None),
], top=3.5)
pageno(s, 2)

# ---------- 3 問題 ----------
s = slide(); kicker(s, "① PROBLEM"); title(s, "台股投資與流動性的結構性痛點")
bullets(s, [
    ("僅限盤中", "09:00–13:30 才能交易,夜間/假日無法即時反應國際盤。"),
    ("單位門檻", "一張 1,000 股門檻高;零股流動性低、撮合慢。"),
    ("跨境困難", "外資開戶/換匯/稅務繁瑣,進場成本高。"),
    ("估值與擔保不透明", "散戶難即時驗證公允價值與背後擔保。"),
    ("無法進入 DeFi", "持股不能直接質押借貸或組合應用。"),
    ("一句話", "台股是高品質資產,卻被「時段 × 單位 × 國界 × 結算」鎖住流動性。"),
], top=2.55, size=16, gap=10)
pageno(s, 3)

# ---------- 4 解決方案 ----------
s = slide(); kicker(s, "② SOLUTION"); title(s, "FormosaX 解決方案(全部已上鏈)")
bullets(s, [
    ("代幣化台股 (ERC-20)", "每檔台股 1:1 對應代幣(dTSMC…),工廠量產、可碎片化至 0.01 股。"),
    ("預言機即時估值", "Chainlink AggregatorV3 風格,餵入 TWSE 價,依即時價鑄造/贖回。"),
    ("AI 輔助估值與風險（題目七核心）", "回歸模型預測公允價 + 風險分數,動態調整借貸 LTV/清算線。"),
    ("受監管 TWD 穩定幣", "對標金管會草案:100% 儲備證明、揭露、贖回、凍結/暫停、角色分權。"),
    ("DeFi 應用", "質押台股代幣借 TWD、代幣化新股認購(超額 pro-rata 配額)。"),
], top=2.55, size=16, gap=11)
pageno(s, 4)

# ---------- 5 系統架構 ----------
s = slide(); kicker(s, "ARCHITECTURE · 技術建議對照"); title(s, "系統架構")
box = s.shapes.add_shape(1, Inches(0.9), Inches(2.55), Inches(11.5), Inches(2.7))
box.fill.solid(); box.fill.fore_color.rgb = CARD; box.line.color.rgb = LINE
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.18)
arch = ("使用者(MetaMask / wagmi + RainbowKit) → React+Vite+shadcn → rwa0607.vercel.app\n"
        "──────────── 鏈上(Ethereum Sepolia,可遷 L2)────────────\n"
        "RegulatedTWD(穩定幣+儲備證明+角色/凍結/暫停)\n"
        "StockFactory ─量產─► StockToken ×12(代幣化台股)\n"
        "PriceOracle(Chainlink 風格) · LendingPool(借貸/清算) · StockIPO(認購)\n"
        "▲ 鏈下 Feeder(TWSE OpenAPI 餵價)   ▲ AI 估值與風險引擎(Python 回歸)")
setp(tf.paragraphs[0], arch, 13, RGBColor(0xBF, 0xE9, 0xDC), font=MONO)
bullets(s, [
    ("區塊鏈 / 標準", "Sepolia 可平移 Arbitrum/Base;ERC-20 碎片化+DeFi,ERC-721 預留不動產類。"),
    ("Oracle", "自建 Chainlink 風格;路線圖以 Chainlink Functions 去中心化餵價並把 AI 輸出上鏈。"),
], top=5.45, size=14, gap=6)
pageno(s, 5)

# ---------- 6 AI 估值 ----------
s = slide(); kicker(s, "AI VALUATION & RISK · 題目七核心"); title(s, "AI 輔助估值與風險模型")
cards(s, [
    ("① 公允估值(回歸)", "歷史價量+技術指標+大盤連動+財報因子 → 線性回歸/隨機森林預測 T+1 公允價區間,標示高/低估。", JADE),
    ("② 風險評分", "波動度、流動性、抵押集中度 → 0–100 風險分數,餵給 LendingPool 動態調整 LTV/清算線。", GOLD),
    ("③ 動態估值 Dashboard", "市價 vs AI 公允價偏離、風險熱力圖、抵押健康度、儲備覆蓋率,讀自鏈上 view。", None),
], top=2.7)
bullets(s, [
    "落地路徑:鏈下 Python 訓練/推論 → 簽章 → Chainlink Functions 把公允價/風險分數寫上鏈供合約引用。",
], top=5.0, size=14)
pageno(s, 6)

# ---------- 7 市場 ----------
s = slide(); kicker(s, "③ TARGET CUSTOMERS / MARKET"); title(s, "目標客戶與市場規模")
cards(s, [
    ("加密散戶 / DeFi", "用穩定幣 24/7 配置台股、賺鏈上收益。", None),
    ("海外華人 / 外資", "低門檻參與台股,免開券商戶、無國界。", None),
    ("券商 / 資管機構", "代幣化發行、做市與新商品通路。", None),
], top=2.6)
bullets(s, [
    ("TAM", "台股總市值約 NT$70 兆+;全球 RWA 代幣化 BCG 預估 2030 達 ~US$16 兆。"),
    ("SAM / SOM", "可代幣化精選台股+跨境/DeFi 客群;初期 30–50 檔藍籌「報價全市場、可交易精選」。"),
    ("趨勢", "RWA 鏈上規模 2023→2025 高速成長,台股資料公開、估值明確,最適合亞洲 RWA 示範。"),
], top=4.7, size=14, gap=7)
pageno(s, 7)

# ---------- 8 商業模式 ----------
s = slide(); kicker(s, "④ BUSINESS MODEL"); title(s, "商業模式 / 獲利 / 預期效益")
table(s, ["收入來源", "說明"], [
    ["交易手續費", "鑄造/贖回微費(Demo 暫免,商轉可開)"],
    ["AUM 管理費", "對代幣化資產收年化管理費"],
    ["借貸利差", "LendingPool 借/存利差 + 清算手續費"],
    ["IPO 承銷費", "StockIPO 代幣化新股認購承銷費"],
    ["穩定幣浮存", "TWD 儲備之利息收益(對標真實穩定幣商模)"],
], top=2.6, fsize=14)
bullets(s, ["效益:投資人(零股化/24-7/可進 DeFi)、發行方(新通路+流動性+可程式化合規)、平台(多元收入+網路效應)。"], top=5.9, size=13)
pageno(s, 8)

# ---------- 9 競爭優勢 ----------
s = slide(); kicker(s, "⑤ COMPETITIVE ADVANTAGE"); title(s, "為什麼比傳統金融好")
table(s, ["構面", "傳統券商", "中心化交易所", "FormosaX"], [
    ["交易時間", "盤中", "24/7", "24/7"],
    ["單位", "整股/零股", "受限", "可碎片化 0.01 股"],
    ["結算", "T+2", "內部帳", "鏈上即時"],
    ["估值透明", "低", "中", "Oracle + AI,鏈上可驗"],
    ["擔保證明", "不可即時", "託管黑箱", "鏈上 1:1 Proof of Reserve"],
    ["可組合性", "無", "有限", "質押/借貸/認購/DeFi"],
    ["合規", "高", "不一", "穩定幣 + 角色/凍結/暫停"],
], top=2.5, fsize=13)
pageno(s, 9)

# ---------- 10 Lean Canvas ----------
s = slide(); kicker(s, "LEAN CANVAS"); title(s, "精實畫布")
table(s, None, [
    ["① 客戶區隔", "加密散戶、海外華人/外資、券商/資管", "② 價值主張", "24/7 碎片化台股+透明估值+DeFi+合規穩定幣"],
    ["③ 通路", "Web DApp(已上線)、錢包、券商 API", "④ 客戶關係", "自助 DApp、儲備/估值儀表、社群"],
    ["⑤ 收入流", "交易費、AUM、借貸利差、IPO、浮存", "⑥ 關鍵資源", "智能合約、預言機、AI 模型、實股託管"],
    ["⑦ 關鍵活動", "餵價/估值、合約維運、做市、合規", "⑧ 關鍵夥伴", "券商/保管行、Chainlink、做市商、審計"],
    ["⑨ 成本結構", "gas/L2、雲端 AI、審計、法遵、開發維運", "", ""],
], top=2.5, fsize=12)
pageno(s, 10)

# ---------- 11 Tokenomics ----------
s = slide(); kicker(s, "TOKENOMICS CANVAS"); title(s, "代幣經濟")
bullets(s, [
    ("名稱/定位", "dXXXX(dTSMC)=代幣化台股,1 代幣↔1 股經濟權利;TWD=平台計價穩定幣。"),
    ("背後資產", "台灣上市股票,價格來源 TWSE;商轉以實股 1:1 託管擔保。"),
    ("Mint / Burn", "存 TWD 依即時價 mint;redeem 反向燒毀;穩定幣 mintTWD/requestRedemption。"),
    ("發行量", "無固定上限,供給隨 TWD 抵押鑄造彈性增減,恆等於儲備擔保量。"),
    ("角色", "admin/issuer、ATTESTOR(儲備簽證)、COMPLIANCE(凍結/暫停)、工廠 owner。"),
    ("費用 / 二級市場", "Demo 免費;商轉開交易費/AUM/利差;未來上 DEX 做市 + 借貸釋放流動性。"),
], top=2.55, size=15, gap=9)
pageno(s, 11)

# ---------- 12 Use Case ----------
s = slide(); kicker(s, "USE CASE CANVAS"); title(s, "使用情境畫布")
cards(s, [
    ("目標使用者", "機構(發行/做市)、個人(配置/賺息)、零售(零股/跨境)。", JADE),
    ("場景 / 痛點解方", "夜間配置、零股碎片化、抵押借 TWD、IPO 認購、即時驗證估值。", JADE),
    ("鏈上 / 鏈下流程", "鏈下 TWSE 餵價+AI 估值 → 鏈上鑄造/贖回/借貸/認購。", JADE),
    ("法規 / 合規挑戰", "證券型代幣定性、KYC/AML、穩定幣儲備與贖回、跨境稅務。", GOLD),
    ("競品 / 差異化", "vs 券商(24/7+碎片化)、vs CEX(可驗擔保)、vs RWA(Oracle+AI+DeFi)。", GOLD),
    ("風險管理機制", "凍結地址、暫停、黑名單、預言機過期價檢查、健康因子/清算。", GOLD),
], top=2.6)
pageno(s, 12)

# ---------- 13 Role Map ----------
s = slide(); kicker(s, "SMART CONTRACT ROLE MAP"); title(s, "合約角色與權限")
table(s, ["角色", "權限"], [
    ["DEFAULT_ADMIN / issuer", "發行、設定揭露、issue 鑄造"],
    ["ATTESTOR", "attestReserves 簽署儲備證明(時間戳 + 文件雜湊)"],
    ["COMPLIANCE", "setFrozen 凍結地址、pause / unpause 全面暫停"],
    ["owner(工廠/借貸/IPO)", "上架代幣、設風險參數、建立 IPO 認購案"],
    ["feeder", "updatePrices 餵價"],
], top=2.5, fsize=13)
bullets(s, [
    "合約:RegulatedTWD·PriceOracle·StockFactory·StockToken×12·LendingPool·StockIPO(皆已驗證)。",
    "轉移:_update 攔截黑名單/凍結+whenNotPaused;現不可升級(易審計),路線圖 UUPS+多簽;permit 路線圖支援;事件全上鏈。",
], top=5.55, size=12, gap=6)
pageno(s, 13)

# ---------- 14 實作亮點 ----------
s = slide(); kicker(s, "LIVE DEMO · 已真實上鏈"); title(s, "實作亮點　rwa0607.vercel.app")
bullets(s, [
    ("模擬資產上鏈", "12 檔台股代幣化,可碎片化買賣。"),
    ("動態估值 Dashboard", "即時市價、儲備覆蓋率、借貸健康因子,讀自鏈上 view。"),
    ("完整 DeFi", "領 TWD→買賣→質押借 TWD→認購 IPO→穩定幣贖回,每筆附 Etherscan 連結。"),
], top=2.5, size=15, gap=8)
box = s.shapes.add_shape(1, Inches(0.9), Inches(4.5), Inches(11.5), Inches(2.3))
box.fill.solid(); box.fill.fore_color.rgb = CARD; box.line.color.rgb = LINE
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.14)
setp(tf.paragraphs[0], "合約地址(Sepolia · Etherscan 已驗證綠勾)", 12, MUT)
for nm, ad, col in [
    ("RegulatedTWD ", "0x0ec97A83E96C7e61b4eDD6FB5Bf33ddEaF63df67", JADE),
    ("PriceOracle  ", "0xF1198ab9A92E21E60bB90003eDEC0887aecf2871", JADE),
    ("LendingPool  ", "0xBebab317ec999ced702dE87C654868a9ede9ab82", JADE),
    ("StockIPO     ", "0x925BE0f709e4c9b38578142E9B94AFbb71aBB990", JADE),
    ("dTSMC(台積電)", "0x1CB2Ea6f8eB47c24387B015C75D9d1dad800dc53", GOLD),
]:
    p = tf.add_paragraph(); p.text = nm + "  " + ad
    p.runs[0].font.size = Pt(12); p.runs[0].font.name = MONO; p.runs[0].font.color.rgb = col
pageno(s, 14)

# ---------- 15 流動性 ----------
s = slide(); kicker(s, "LIQUIDITY IMPACT"); title(s, "流動性提升效果")
table(s, ["指標", "傳統台股", "FormosaX 代幣化"], [
    ["交易時間", "~4.5 小時/日", "24 小時 × 7 天"],
    ["最小單位", "1 股(零股慢)", "0.01 股(完全碎片化)"],
    ["結算", "T+2", "即時(鏈上原子結算)"],
    ["跨境參與", "高門檻", "錢包即可,無國界"],
    ["可組合性", "無", "質押/借貸/認購/DeFi"],
    ["擔保驗證", "不可即時", "鏈上 1:1 Proof of Reserve"],
], top=2.5, fsize=13)
bullets(s, ["交易可及時間 +400%、單位門檻 ↓99%、新增抵押釋放流動性管道 → 有效流動性顯著提升。"], top=6.0, size=13)
pageno(s, 15)

# ---------- 16 二級市場 Uniswap ----------
s = slide(); kicker(s, "SECONDARY MARKET · UNISWAP"); title(s, "二級市場:Uniswap 流動性池")
cards(s, [
    ("站內 AMM 交換", "前端直接打 Uniswap v3 SwapRouter,吃同一池流動性;QuoterV2 報價、exactInputSingle 交換。", JADE),
    ("外部 Uniswap 可換", "同一池在 app.uniswap.org(Sepolia)用 token 位址匯入即可交易。", JADE),
    ("一/二級市場套利收斂", "預言機價(一級)錨定公允價、AMM(二級)做 24/7 價格發現,偏離由套利者拉回。", GOLD),
], top=2.6)
tf = tb(s, 0.9, 4.8, 11.6, 0.6)
setp(tf.paragraphs[0], "dTSMC/TWD v3 池(0.3%)  0xe69aF036F8290Ed0371f79dE8E1060211e9f5353  ·  池價 = 預言機價(已驗證)", 13, JADE, font=MONO)
bullets(s, ["把「流動性提升」做成真的可交易二級市場 —— 不只是 mint/redeem,而是 24/7 AMM 自由買賣 + 真實價格發現。"], top=5.6, size=14)
pageno(s, 16)

# ---------- 17 結語 ----------
s = slide(); kicker(s, "ROADMAP · CLOSING"); title(s, "路線圖與結語")
cards(s, [
    ("Phase 1 · 已完成", "台股代幣化+預言機+穩定幣+借貸+IPO,部署 Sepolia、Etherscan 驗證、上線 Vercel。", JADE),
    ("Phase 2", "AI 估值/風險上鏈(Chainlink Functions)、動態 LTV、遷 L2 降 gas。", None),
    ("Phase 3", "實股託管+KYC/合規、上 DEX 做市、擴展不動產/債權(ERC-721)。", GOLD),
], top=2.6)
tf = tb(s, 0.9, 4.8, 11.6, 1.0)
setp(tf.paragraphs[0], "FormosaX 以台股示範「透明估值(Oracle+AI)× 鏈上擔保 × 可組合 DeFi × 可程式化合規」的完整 RWA 堆疊,且已真實部署、可驗證。", 15, MUT)
tf = tb(s, 0.9, 5.9, 11.6, 1.0)
setp(tf.paragraphs[0], "Thanks for watching — Q&A", 32, JADE, bold=True)
pageno(s, 17)

out = r"C:\Users\lab643\Desktop\rwa report\FormosaX_RWA_Deck.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
